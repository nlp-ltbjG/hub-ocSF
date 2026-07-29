"""
Harness 核心引擎 - 渐进式披露架构

基于 Harness Engineering 设计理念：
- 常驻层: Skill 索引始终加载
- 触发层: 匹配条件后按需加载完整 Skill
- 执行层: Skill 执行期间完整驻留
- LLM: deepseek-v4-flash 智能推理
"""

import os
import re
import time
from pathlib import Path
from typing import Dict, List, Optional

from llm_client import LLMClient, create_llm_client


class ProgressiveDisclosure:
    """渐进式披露管理器"""

    def __init__(self, config_dir: str):
        self.config_dir = Path(config_dir)
        self.skill_index = self._load_skill_index()
        self.loaded_skills: Dict[str, dict] = {}

    def _load_skill_index(self) -> List[dict]:
        """加载常驻层: Skill 索引"""
        index_path = self.config_dir / "skill_index.md"
        if not index_path.exists():
            return []

        skills = []
        content = index_path.read_text(encoding="utf-8")
        current_skill = None

        for line in content.split("\n"):
            line = line.strip()
            if line.startswith("## ") and line.strip() == "## 可用 Skills":
                continue
            if line.startswith("- `"):
                if current_skill:
                    skills.append(current_skill)
                current_skill = {"triggers": [], "description": ""}
                match = re.match(r"- `([^`]+)`\s*—\s*(.*)", line)
                if match:
                    current_skill["name"] = match.group(1)
                    current_skill["description"] = match.group(2)
            elif current_skill and "触发词:" in line:
                triggers_str = line.split("触发词:")[1].strip()
                current_skill["triggers"] = [t.strip() for t in triggers_str.split("|")]
            elif current_skill and "描述:" in line:
                current_skill["description"] = line.split("描述:")[1].strip()

        if current_skill:
            skills.append(current_skill)

        return skills

    def match_skill(self, user_input: str) -> Optional[dict]:
        """根据用户输入匹配 Skill"""
        user_input_lower = user_input.lower()
        for skill in self.skill_index:
            for trigger in skill.get("triggers", []):
                if trigger.lower() in user_input_lower:
                    return skill
        return None

    def load_skill(self, skill_name: str) -> Optional[str]:
        """按需加载完整 Skill 定义"""
        if skill_name in self.loaded_skills:
            return self.loaded_skills[skill_name].get("full_content", "")

        skill_path = self.config_dir.parent / "skills" / skill_name / "SKILL.md"
        if not skill_path.exists():
            return None

        content = skill_path.read_text(encoding="utf-8")
        self.loaded_skills[skill_name] = {
            "full_content": content,
            "loaded_at": time.time(),
        }
        return content

    def release_skill(self, skill_name: str):
        """释放已加载的 Skill"""
        self.loaded_skills.pop(skill_name, None)

    def get_index_summary(self) -> str:
        """获取常驻层摘要"""
        lines = ["# Skill 索引（常驻层）\n"]
        for skill in self.skill_index:
            triggers = ", ".join(skill.get("triggers", [])[:3])
            lines.append(f"- `{skill['name']}` — {skill.get('description', '')}")
            lines.append(f"  触发词: {triggers}")
        return "\n".join(lines)


class MemorySystem:
    """四层记忆模型"""

    def __init__(self, config_dir: str):
        self.config_dir = Path(config_dir)
        self.working_memory: List[dict] = []
        self.long_term_memory = self._load_long_term_memory()

    def _load_long_term_memory(self) -> dict:
        """加载 L3 长期记忆"""
        memory = {}
        mem_path = self.config_dir / "MEMORY.md"
        if mem_path.exists():
            content = mem_path.read_text(encoding="utf-8")
            memory["raw"] = content
            memory["sections"] = self._parse_memory_sections(content)
        return memory

    def _parse_memory_sections(self, content: str) -> Dict[str, str]:
        """解析 MEMORY.md 的各个段落"""
        sections = {}
        current_section = None
        for line in content.split("\n"):
            if line.startswith("## "):
                current_section = line[3:].strip()
                sections[current_section] = []
            elif current_section and line.strip():
                sections[current_section].append(line)
        return sections

    def add_to_working(self, role: str, content: str):
        """添加到 L1 工作记忆"""
        self.working_memory.append(
            {"role": role, "content": content, "timestamp": time.time()}
        )

    def get_context(self) -> str:
        """组装上下文"""
        parts = []

        if self.long_term_memory.get("sections"):
            for section, lines in self.long_term_memory["sections"].items():
                if section in ["Skill 索引（常驻层）", "用户偏好", "技术栈"]:
                    parts.append(f"## {section}\n" + "\n".join(lines))

        recent = self.working_memory[-10:]
        if recent:
            parts.append("## 最近对话")
            for msg in recent:
                parts.append(f"[{msg['role']}]: {msg['content']}")

        return "\n\n".join(parts)

    def flush(self):
        """记忆刷新"""
        if self.working_memory:
            print(f"[Memory Flush] 刷新 {len(self.working_memory)} 条工作记忆")
            self.working_memory = []


class Gateway:
    """Fat Gateway: 消息网关与会话管理"""

    def __init__(self):
        self.sessions: Dict[str, dict] = {}
        self.lane_queues: Dict[str, List[dict]] = {}

    def create_session(self, session_id: str, user_id: str = "default") -> dict:
        """创建新会话"""
        session = {
            "sessionId": session_id,
            "userId": user_id,
            "channel": "harness",
            "createdAt": time.time(),
            "isRunning": False,
            "hasError": False,
            "retryCount": 0,
        }
        self.sessions[session_id] = session
        self.lane_queues[session_id] = []
        return session

    def submit_message(self, session_id: str, content: str) -> dict:
        """提交消息到 Lane 队列"""
        if session_id not in self.sessions:
            self.create_session(session_id)

        lane = self.lane_queues[session_id]
        message = {"content": content, "timestamp": time.time(), "processed": False}
        lane.append(message)
        print(f"[Gateway] 消息入队 (session={session_id}, queue_size={len(lane)})")
        return message

    def get_next_message(self, session_id: str) -> Optional[dict]:
        """获取队列中下一条消息"""
        if session_id not in self.lane_queues:
            return None

        lane = self.lane_queues[session_id]
        for msg in lane:
            if not msg["processed"]:
                return msg
        return None

    def mark_processed(self, session_id: str, message: dict):
        """标记消息为已处理"""
        message["processed"] = True

    def get_session_status(self, session_id: str) -> dict:
        """获取会话状态"""
        if session_id not in self.sessions:
            return {"error": "Session not found"}
        session = self.sessions[session_id]
        lane = self.lane_queues.get(session_id, [])
        pending = [m for m in lane if not m["processed"]]
        return {
            "sessionId": session_id,
            "isRunning": session["isRunning"],
            "pendingCount": len(pending),
            "retryCount": session["retryCount"],
        }


class AgentNode:
    """Agent 执行引擎: ReAct 循环"""

    def __init__(
        self,
        gateway: Gateway,
        memory: MemorySystem,
        disclosure: ProgressiveDisclosure,
        llm: LLMClient = None,
    ):
        self.gateway = gateway
        self.memory = memory
        self.disclosure = disclosure
        self.llm = llm
        self.max_turns = 10
        self.tools = self._register_tools()
        self.base_dir = Path(__file__).resolve().parent.parent
        self.result_dir = self.base_dir / "result"
        self.result_dir.mkdir(parents=True, exist_ok=True)

    def _register_tools(self) -> dict:
        """注册可用工具"""
        return {
            "ppt_extractor": self._tool_ppt_extractor,
            "knowledge_summarizer": self._tool_knowledge_summarizer,
            "web_searcher": self._tool_web_searcher,
            "markdown_writer": self._tool_markdown_writer,
        }

    def run(self, session_id: str, user_input: str) -> dict:
        """执行 ReAct 循环"""
        print(f"\n[AgentNode] 开始执行 (session={session_id})")
        print(f"[AgentNode] 用户输入: {user_input[:100]}...\n")

        self.memory.add_to_working("user", user_input)

        matched_skill = self.disclosure.match_skill(user_input)
        if not matched_skill:
            return self._handle_no_match(user_input)

        skill_content = self.disclosure.load_skill(matched_skill["name"])
        if not skill_content:
            return {
                "status": "error",
                "message": f"Skill {matched_skill['name']} 加载失败",
            }

        print(f"[Progressive Disclosure] 匹配 Skill: {matched_skill['name']}")
        print(f"[Progressive Disclosure] 按需加载完整定义 ({len(skill_content)} chars)")

        result = self._execute_skill(matched_skill["name"], skill_content, user_input)

        self.disclosure.release_skill(matched_skill["name"])
        print(f"[Progressive Disclosure] 释放 Skill {matched_skill['name']}")

        return result

    def _execute_skill(
        self, skill_name: str, skill_content: str, user_input: str
    ) -> dict:
        """执行 Skill 流程，共享 execution_context"""
        steps = self._parse_skill_steps(skill_content)
        results = []
        execution_context = {
            "user_input": user_input,
            "ppt_files": [],
            "extracted_texts": {},
            "summary": {},
            "supplementary_topics": [],
            "search_strategies": [],
            "output_path": None,
            "topic": "",
        }

        for i, step in enumerate(steps):
            print(f"\n[ReAct] 第 {i + 1} 步: {step['name']}")
            print(f"[Reason] {step.get('reasoning', '')[:100]}...")

            tool_name = step.get("tool", "")
            if tool_name in self.tools:
                tool_result = self.tools[tool_name](execution_context, step)
                results.append(
                    {"step": step["name"], "tool": tool_name, "result": tool_result}
                )
                print(f"[Act] 调用 {tool_name}")
                print(f"[Observation] {str(tool_result)[:200]}...")

                if (
                    isinstance(tool_result, dict)
                    and tool_result.get("status") == "error"
                ):
                    print(f"\n[ReAct] ❌ 工具执行出错，中止后续步骤")
                    return {
                        "status": "error",
                        "skill": skill_name,
                        "steps_completed": len(results),
                        "results": results,
                        "message": tool_result.get("message", "Unknown error"),
                        "error_step": step["name"],
                    }
            else:
                results.append(
                    {"step": step["name"], "tool": None, "result": "Skipped"}
                )

        return {
            "status": "success",
            "skill": skill_name,
            "steps_completed": len(results),
            "results": results,
            "output_path": execution_context.get("output_path"),
            "summary": execution_context.get("summary"),
        }

    def _parse_skill_steps(self, skill_content: str) -> List[dict]:
        """解析 Skill 执行步骤"""
        steps = [
            {
                "name": "PPT 文字提取",
                "reasoning": "从 PPT 文件中提取原始文字内容",
                "tool": "ppt_extractor",
            },
            {
                "name": "知识总结",
                "reasoning": "从提取的内容中提炼核心概念和架构",
                "tool": "knowledge_summarizer",
            },
            {
                "name": "补充缺失知识",
                "reasoning": "为初学者识别需补充的前置知识",
                "tool": "web_searcher",
            },
            {
                "name": "生成 Markdown",
                "reasoning": "将总结内容输出为结构化文档",
                "tool": "markdown_writer",
            },
        ]
        return steps

    def _tool_ppt_extractor(self, ctx: dict, step: dict) -> dict:
        """工具: PPT 文字提取"""
        try:
            from pptx import Presentation
        except ImportError:
            return {
                "status": "error",
                "message": "python-pptx 未安装，请执行: pip install python-pptx",
            }

        ppt_paths = self._extract_ppt_paths(ctx["user_input"])
        if not ppt_paths:
            return {"status": "error", "message": "未找到 PPT 文件路径"}

        ctx["ppt_files"] = ppt_paths
        extracted = {}
        for path in ppt_paths:
            try:
                prs = Presentation(path)
                slides = []
                for i, slide in enumerate(prs.slides):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text = para.text.strip()
                                if text:
                                    texts.append(text)
                    if texts:
                        slides.append({"page": i + 1, "content": texts})
                extracted[path] = slides
                print(f"  [PPT Extractor] {Path(path).name}: {len(slides)} 页提取成功")
            except Exception as e:
                extracted[path] = {"error": str(e)}

        ctx["extracted_texts"] = extracted
        return {
            "status": "success",
            "files_processed": len(ppt_paths),
            "total_slides": sum(
                len(v) for v in extracted.values() if isinstance(v, list)
            ),
        }

    def _tool_knowledge_summarizer(self, ctx: dict, step: dict) -> dict:
        """工具: 知识总结 - 优先使用 LLM，降级为智能规则提取"""
        extracted = ctx.get("extracted_texts", {})
        if not extracted:
            return {"status": "error", "message": "没有可总结的内容"}

        structured_slides = []
        all_texts = []
        for filepath, slides in extracted.items():
            if isinstance(slides, list):
                for slide in slides:
                    structured_slides.append(slide)
                    all_texts.extend(slide.get("content", []))

        combined_text = "\n".join(all_texts)

        if self.llm and self.llm.is_available:
            try:
                print(
                    "  [Knowledge Summarizer] 使用 LLM (deepseek-v4-flash) 进行智能分析..."
                )
                summary = self.llm.summarize_ppt(combined_text)
                if "error" not in summary and summary.get("core_concepts"):
                    topic = summary.get("topic", "") or self._infer_topic(all_texts)
                    summary["topic"] = topic
                    ctx["summary"] = summary
                    ctx["topic"] = topic
                    concepts_count = len(summary["core_concepts"])
                    print(
                        f"  [Knowledge Summarizer] ✅ LLM 分析成功: {concepts_count} 个核心概念"
                    )
                    print(f"  [Knowledge Summarizer] 🎯 主题: {topic}")
                    return {
                        "status": "success",
                        "core_concepts_found": concepts_count,
                        "topic": topic,
                        "llm_used": True,
                    }
                else:
                    print("  [Knowledge Summarizer] ⚠️ LLM 输出异常，降级为智能规则提取")
            except Exception as e:
                print(
                    f"  [Knowledge Summarizer] ⚠️ LLM 调用失败 ({e})，降级为智能规则提取"
                )

        print("  [Knowledge Summarizer] 🔧 使用智能规则提取模式")
        summary = self._analyze_slides(structured_slides, all_texts)
        ctx["summary"] = summary
        topic = summary.get("topic", "") or self._infer_topic(all_texts)
        summary["topic"] = topic
        ctx["topic"] = topic
        print(f"  [Knowledge Summarizer] 🎯 主题: {topic}")
        print(
            f"  [Knowledge Summarizer] 📊 提取: {len(summary.get('core_concepts', []))} 概念, "
            f"{len(summary.get('architectures', []))} 架构, "
            f"{len(summary.get('comparisons', []))} 对比"
        )

        return {
            "status": "success",
            "core_concepts_found": len(summary.get("core_concepts", [])),
            "topic": topic,
            "llm_used": False,
        }

    def _analyze_slides(self, slides: list, all_texts: list) -> dict:
        """基于幻灯片结构的智能分析 - 降级模式的核心"""
        concepts = []
        architectures = []
        comparisons = []
        best_practices = []
        applications = []

        slide_titles = []
        for slide in slides:
            content = slide.get("content", [])
            if content:
                title = self._extract_slide_title(content)
                if title:
                    slide_titles.append(title)

        seen_concepts = set()

        for slide in slides:
            content = slide.get("content", [])
            if not content:
                continue

            title = self._extract_slide_title(content)
            body = content[1:] if title else content

            if title and len(title) > 3:
                concepts.append({"name": title, "description": ""})
                seen_concepts.add(title.strip().lower())

            for line in body:
                clean = self._clean_line(line)
                if not clean or len(clean) < 4:
                    continue
                clean_lower = clean.lower()

                if clean_lower in seen_concepts:
                    continue

                if self._is_architecture_line(clean):
                    architectures.append({"name": clean, "description": ""})
                    seen_concepts.add(clean_lower)
                    continue

                if self._is_comparison_line(clean):
                    comp = self._parse_comparison(clean)
                    if comp:
                        comparisons.append(comp)
                    seen_concepts.add(clean_lower)
                    continue

                if self._is_practice_line(clean):
                    best_practices.append({"title": clean, "description": ""})
                    seen_concepts.add(clean_lower)
                    continue

                if self._is_application_line(clean):
                    applications.append({"scenario": clean, "description": ""})
                    seen_concepts.add(clean_lower)
                    continue

                if self._is_concept_line(clean):
                    concepts.append({"name": clean, "description": ""})
                    seen_concepts.add(clean_lower)

        if len(concepts) > 20:
            concepts = concepts[:20]

        concepts = self._filter_metadata_concepts(concepts)

        if not concepts:
            for line in all_texts[:30]:
                clean = self._clean_line(line)
                if clean and len(clean) > 4:
                    concepts.append({"name": clean[:60], "description": ""})
                if len(concepts) >= 10:
                    break

        topic = self._derive_topic(slide_titles, concepts, all_texts)

        supplementary = self._derive_supplementary_topics(
            concepts, architectures, comparisons
        )

        return {
            "core_concepts": concepts,
            "architectures": architectures,
            "comparisons": comparisons,
            "best_practices": best_practices,
            "applications": applications,
            "supplementary_topics": supplementary,
            "topic": topic,
        }

    def _filter_metadata_concepts(self, concepts: list) -> list:
        """过滤元数据类的概念条目（版权、来源、宣传语等）"""
        metadata_patterns = [
            r"出品",
            r"出品方",
            r"版权",
            r"盗版必究",
            r"GitHub\s*Star",
            r"万\+",
            r"www\.",
            r"http",
            r"扫码",
            r"关注",
            r"公众号",
            r"微信",
            r"进度\d+",
            r"完成度",
        ]
        filtered = []
        for concept in concepts:
            name = (
                concept.get("name", "") if isinstance(concept, dict) else str(concept)
            )
            is_metadata = False
            for pattern in metadata_patterns:
                if re.search(pattern, name, re.IGNORECASE):
                    is_metadata = True
                    break
            if not is_metadata and len(name) > 3:
                filtered.append(concept)
        return filtered

    def _clean_line(self, line: str) -> str:
        """清理文本行：去除标记符号和多余空白"""
        clean = line.strip()
        clean = re.sub(r"^[•·\-\*–—\+\d\.\)、]+", "", clean)
        clean = re.sub(r"\s+", " ", clean)
        return clean.strip()

    def _extract_slide_title(self, content: list) -> str:
        """从幻灯片内容中提取标题（通常是第一个有效且较短的文本行）"""
        for text in content:
            clean = self._clean_line(text)
            if 3 < len(clean) <= 50:
                return clean
            if len(clean) > 50:
                return ""
        return ""

    def _is_concept_line(self, text: str) -> bool:
        """判断是否为概念/术语行"""
        concept_markers = [
            "是",
            "是一种",
            "是一个",
            "称为",
            "叫做",
            "指",
            "即",
            "核心",
            "关键",
            "重要",
            "主要",
            "基本",
            "Engine",
            "Framework",
            "Library",
            "Platform",
            "System",
            "Module",
            "Component",
            "Service",
            "Model",
            "Algorithm",
            "Protocol",
            "Interface",
            "架构",
            "框架",
            "模型",
            "协议",
            "接口",
        ]
        if len(text) < 3 or len(text) > 60:
            return False
        for marker in concept_markers:
            if marker.lower() in text.lower():
                return True
        if text.endswith(("。", "！", "？", "：", ":")):
            return False
        if len(text) <= 30 and not text.startswith(("因为", "所以", "但是", "如果")):
            return True
        return False

    def _is_architecture_line(self, text: str) -> bool:
        """判断是否为架构/结构行"""
        arch_keywords = [
            "架构",
            "模型",
            "分层",
            "组件",
            "模块",
            "流程",
            "管道",
            "架构图",
            "结构图",
            "分层",
            "Stack",
            "Layer",
            "Pipeline",
            "组件",
            "节点",
            "适配器",
            "网关",
            "路由器",
            "三层",
            "四层",
            "五层",
            "六层",
        ]
        for kw in arch_keywords:
            if kw.lower() in text.lower():
                return True
        if re.match(r"^[一二三四五六七八九十]+[层级章节]", text):
            return True
        if re.match(r"^\d+[\.、)）]", text) and len(text) > 5:
            return True
        return False

    def _is_comparison_line(self, text: str) -> bool:
        """判断是否为对比行"""
        comp_keywords = [
            "对比",
            "区别",
            "vs",
            "VS",
            "不同于",
            "差异",
            "比较",
            "vs",
            "相比",
            "vs.",
            "相较于",
            "与...对比",
            "Platform vs",
            "Framework vs",
        ]
        for kw in comp_keywords:
            if kw.lower() in text.lower():
                return True
        if " vs " in text.lower() and len(text) > 5:
            return True
        return False

    def _parse_comparison(self, text: str) -> dict:
        """解析对比行，提取 A vs B 结构"""
        vs_patterns = [
            r"(.+?)\s*vs\s*(.+)",
            r"(.+?)\s*VS\s*(.+)",
            r"(.+?)对比(.+)",
            r"(.+?)不同于(.+)",
            r"(.+?)区别(.+)",
        ]
        for pattern in vs_patterns:
            match = re.match(pattern, text, re.IGNORECASE)
            if match:
                a = match.group(1).strip()
                b = match.group(2).strip()
                return {"a": a, "b": b, "difference": text}
        return {"a": text, "b": "", "difference": ""}

    def _is_practice_line(self, text: str) -> bool:
        """判断是否为最佳实践/原则行"""
        practice_keywords = [
            "原则",
            "最佳实践",
            "设计原则",
            "应该",
            "必须",
            "建议",
            "实践",
            "规范",
            "约定",
            "模式",
            "需要",
            "应当",
            "务必",
            "推荐",
        ]
        for kw in practice_keywords:
            if kw in text:
                return True
        if text.startswith(("• ", "- ", "* ")) and len(text) > 10:
            return True
        return False

    def _is_application_line(self, text: str) -> bool:
        """判断是否为应用场景行"""
        app_keywords = [
            "应用",
            "场景",
            "案例",
            "适用",
            "可以用于",
            "使用场景",
            "典型应用",
            "实际案例",
            "适用于",
            "可用于",
            "能够",
        ]
        for kw in app_keywords:
            if kw in text:
                return True
        return False

    def _derive_topic(self, slide_titles: list, concepts: list, all_texts: list) -> str:
        """从幻灯片标题和概念中推断主题"""
        if slide_titles:
            first_title = slide_titles[0]
            if len(first_title) > 3:
                return self._normalize_topic(first_title)

        for concept in concepts:
            name = (
                concept.get("name", "") if isinstance(concept, dict) else str(concept)
            )
            if name and len(name) > 3 and not name.endswith(("。", "！", "？")):
                return self._normalize_topic(name)

        for text in all_texts[:5]:
            clean = self._clean_line(text)
            if 3 < len(clean) < 50:
                return self._normalize_topic(clean)

        return "Knowledge_Summary"

    def _normalize_topic(self, topic: str) -> str:
        """将主题名规范化为可用的标识符"""
        topic = topic.strip()
        if not topic:
            return "Knowledge_Summary"
        if re.match(r"^[\x00-\x7F]+$", topic):
            topic = topic.replace(" ", "_")
        else:
            chinese_chars = "".join(c for c in topic if "\u4e00" <= c <= "\u9fff")
            if chinese_chars:
                pinyin_map = {
                    "系统": "System",
                    "架构": "Architecture",
                    "模型": "Model",
                    "设计": "Design",
                    "开发": "Development",
                    "应用": "Application",
                    "框架": "Framework",
                    "平台": "Platform",
                    "服务": "Service",
                    "工具": "Tool",
                    "技能": "Skill",
                    "记忆": "Memory",
                    "协议": "Protocol",
                    "接口": "Interface",
                    "组件": "Component",
                    "模块": "Module",
                    "安全": "Security",
                    "数据": "Data",
                    "分析": "Analysis",
                    "算法": "Algorithm",
                    "网络": "Network",
                    "前端": "Frontend",
                    "后端": "Backend",
                    "数据库": "Database",
                    "机器学习": "ML",
                    "深度学习": "DL",
                    "人工智能": "AI",
                }
                result = []
                remaining = chinese_chars
                for k, v in sorted(pinyin_map.items(), key=lambda x: -len(x[0])):
                    if k in remaining:
                        result.append(v)
                        remaining = remaining.replace(k, "")
                if result:
                    topic = "_".join(result)
                else:
                    topic = chinese_chars[:10] or "Knowledge_Summary"
            else:
                topic = re.sub(r"[^\w\s&\-]", "", topic)
                topic = topic.replace(" ", "_")
        topic = re.sub(r"_+", "_", topic)
        topic = topic.strip("_&- ")
        return topic if topic else "Knowledge_Summary"

    def _derive_supplementary_topics(
        self, concepts: list, architectures: list, comparisons: list
    ) -> list:
        """推断需要补充的前置知识点"""
        topic_names = set()
        for c in concepts[:5]:
            name = c.get("name", "") if isinstance(c, dict) else str(c)
            topic_names.add(name.lower())

        supplementary = []
        knowledge_domains = [
            ("基础概念", ["基础", "入门", "概念", "介绍", "简介"]),
            ("核心原理", ["原理", "机制", "工作方式", "原理", "底层"]),
            ("实践指南", ["实践", "指南", "教程", "实战", "案例"]),
            ("设计模式", ["模式", "设计", "架构", "结构"]),
            ("性能优化", ["性能", "优化", "效率", "速度"]),
            ("安全基础", ["安全", "权限", "认证", "加密"]),
        ]

        matched = set()
        for domain_name, keywords in knowledge_domains:
            if domain_name in matched:
                continue
            for concept_name in topic_names:
                if any(kw in concept_name for kw in keywords):
                    supplementary.append(domain_name)
                    matched.add(domain_name)
                    break

        if not supplementary:
            primary = ""
            if concepts:
                primary = (
                    concepts[0].get("name", "")
                    if isinstance(concepts[0], dict)
                    else str(concepts[0])
                )
            supplementary = [
                f"{primary}基础概念" if primary else "相关基础概念",
                f"{primary}核心原理" if primary else "核心原理",
                f"{primary}实践指南" if primary else "实践指南",
            ]

        return supplementary[:5]

    def _infer_topic(self, texts: list) -> str:
        """从文本列表推断主题"""
        if not texts:
            return "Knowledge_Summary"
        first_text = self._clean_line(texts[0]) if texts else ""
        if first_text and len(first_text) > 3:
            return self._normalize_topic(first_text)
        return "Knowledge_Summary"

    def _tool_web_searcher(self, ctx: dict, step: dict) -> dict:
        """工具: 补充知识 - 使用 LLM 生成搜索策略，降级为基于概念的智能搜索建议"""
        summary = ctx.get("summary", {})
        topics = summary.get("supplementary_topics", [])
        concepts = summary.get("core_concepts", [])
        topic_name = ctx.get("topic", "")

        search_strategies = []
        llm_used = False

        if self.llm and self.llm.is_available and topics:
            try:
                print(
                    f"  [Web Searcher] 使用 LLM 为 {len(topics)} 个知识点生成搜索策略..."
                )
                search_strategies = self.llm.supplement_knowledge(topics)
                if search_strategies:
                    llm_used = True
                    print(
                        f"  [Web Searcher] ✅ LLM 生成 {len(search_strategies)} 条搜索策略"
                    )
            except Exception as e:
                print(f"  [Web Searcher] ⚠️ LLM 调用失败 ({e})，使用概念驱动策略")

        if not search_strategies:
            search_strategies = self._generate_search_strategies(
                topics, concepts, topic_name
            )
            print(
                f"  [Web Searcher] 🔧 基于概念生成 {len(search_strategies)} 条搜索建议"
            )

        ctx["supplementary_topics"] = topics
        ctx["search_strategies"] = search_strategies

        print(f"  [Web Searcher] 📋 搜索查询列表:")
        for s in search_strategies[:3]:
            print(
                f"    - [{s.get('priority', 'medium')}] {s.get('topic', '')} → {s.get('search_query', '')}"
            )

        return {
            "status": "success",
            "queries_planned": [s.get("search_query", "") for s in search_strategies],
            "topics_to_supplement": topics,
            "search_strategies": search_strategies,
            "llm_used": llm_used,
        }

    def _generate_search_strategies(
        self, topics: list, concepts: list, topic_name: str
    ) -> list:
        """基于概念生成具体的搜索策略"""
        strategies = []

        concept_names = []
        for c in concepts[:8]:
            if isinstance(c, dict):
                name = c.get("name", "")
            else:
                name = str(c)
            if name and len(name) > 2:
                concept_names.append(name)

        seen_queries = set()

        for topic in topics[:5]:
            query = f"{topic} 入门教程 详解"
            if query not in seen_queries:
                seen_queries.add(query)
                priority = "high" if len(strategies) < 2 else "medium"
                strategies.append(
                    {
                        "topic": topic,
                        "search_query": query,
                        "brief_intro": f"理解{topic}的基础概念和核心方法",
                        "priority": priority,
                    }
                )

        for concept in concept_names[:3]:
            query = f"{concept} 基础概念 详解"
            if query not in seen_queries:
                seen_queries.add(query)
                strategies.append(
                    {
                        "topic": concept,
                        "search_query": query,
                        "brief_intro": f"「{concept}」的入门知识",
                        "priority": "medium",
                    }
                )

        if not strategies:
            strategies.append(
                {
                    "topic": topic_name or "相关技术",
                    "search_query": f"{topic_name or '相关技术'} 入门教程 2025",
                    "brief_intro": "系统学习相关技术",
                    "priority": "medium",
                }
            )

        return strategies[:6]

    def _tool_markdown_writer(self, ctx: dict, step: dict) -> dict:
        """工具: Markdown 文档生成"""
        summary = ctx.get("summary", {})
        topic = ctx.get("topic", "Knowledge_Summary")
        supplementary = ctx.get("supplementary_topics", [])
        search_strategies = ctx.get("search_strategies", [])
        ppt_files = ctx.get("ppt_files", [])

        md_content = self._build_markdown(
            summary, topic, supplementary, ppt_files, search_strategies
        )

        safe_name = topic.replace(" ", "_").replace("/", "_")
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        filename = f"{safe_name}_{timestamp}.md"

        output_path = self.result_dir / filename
        output_path.write_text(md_content, encoding="utf-8")

        abs_path = os.path.abspath(str(output_path))
        ctx["output_path"] = abs_path

        print(f"  [Markdown Writer] 文件已生成: {abs_path}")
        print(f"  [Markdown Writer] 文件大小: {len(md_content)} chars")

        return {
            "status": "success",
            "output_path": abs_path,
            "file_size": len(md_content),
            "filename": filename,
        }

    def _build_markdown(
        self,
        summary: dict,
        topic: str,
        supplementary: list,
        ppt_files: list,
        search_strategies: list = None,
    ) -> str:
        """构建 Markdown 内容 - 纯数据驱动，无硬编码"""
        lines = []
        safe_topic = topic.replace("_", " ")
        lines.append(f"# {safe_topic} 知识总结\n")
        lines.append(f"> **生成时间**: {time.strftime('%Y-%m-%d %H:%M:%S')}")
        lines.append(f"> **来源文件**: {', '.join(Path(f).name for f in ppt_files)}")
        lines.append(f"> **分析引擎**: Harness Engine v2.0 (deepseek-v4-flash)")
        lines.append("")

        has_any_content = False

        concepts = summary.get("core_concepts", [])
        if concepts:
            has_any_content = True
            lines.append("## 一、核心概念\n")
            lines.append("| # | 概念 | 说明 |")
            lines.append("|---|------|------|")
            for i, item in enumerate(concepts, 1):
                if isinstance(item, dict):
                    name = item.get("name", item.get("concept", ""))
                    desc = item.get("description", item.get("desc", ""))
                else:
                    name = str(item)
                    desc = ""
                lines.append(f"| {i} | **{name}** | {desc} |")
            lines.append("")

        archs = summary.get("architectures", [])
        if archs:
            has_any_content = True
            lines.append("## 二、架构与模型\n")
            for i, item in enumerate(archs, 1):
                if isinstance(item, dict):
                    name = item.get("name", item.get("architecture", ""))
                    desc = item.get("description", item.get("desc", ""))
                else:
                    name = str(item)
                    desc = ""
                lines.append(f"### 2.{i} {name}\n")
                if desc:
                    lines.append(f"{desc}\n")

        comps = summary.get("comparisons", [])
        if comps:
            has_any_content = True
            lines.append("## 三、对比分析\n")
            lines.append("| 项目A | 项目B | 对比要点 |")
            lines.append("|-------|-------|---------|")
            for item in comps:
                if isinstance(item, dict):
                    a = item.get("a", item.get("item_a", item.get("from", "")))
                    b = item.get("b", item.get("item_b", item.get("to", "")))
                    diff = item.get(
                        "difference", item.get("key_points", item.get("desc", ""))
                    )
                else:
                    a = str(item)
                    b = ""
                    diff = ""
                lines.append(f"| {a} | {b} | {diff} |")
            lines.append("")

        practices = summary.get("best_practices", [])
        if practices:
            has_any_content = True
            lines.append("## 四、最佳实践\n")
            for i, item in enumerate(practices, 1):
                if isinstance(item, dict):
                    title = item.get("title", item.get("name", f"实践{i}"))
                    desc = item.get("description", item.get("desc", ""))
                else:
                    title = str(item)
                    desc = ""
                lines.append(f"- **{title}**")
                if desc:
                    lines.append(f"  - {desc}")
            lines.append("")

        apps = summary.get("applications", [])
        if apps:
            has_any_content = True
            lines.append("## 五、应用场景\n")
            lines.append("| # | 场景 | 说明 |")
            lines.append("|---|------|------|")
            for i, item in enumerate(apps, 1):
                if isinstance(item, dict):
                    name = item.get("scenario", item.get("name", ""))
                    desc = item.get("description", item.get("desc", ""))
                else:
                    name = str(item)
                    desc = ""
                lines.append(f"| {i} | {name} | {desc} |")
            lines.append("")

        if not has_any_content:
            lines.append("> 本 PPT 未提取到结构化内容，建议检查文件格式或内容。\n")

        if search_strategies:
            lines.append("## 六、前置知识补充\n")
            lines.append("以下为理解本主题需要掌握的前置知识点，建议优先学习：\n")
            for i, strategy in enumerate(search_strategies, 1):
                topic_name = strategy.get("topic", f"知识点{i}")
                brief = strategy.get("brief_intro", "")
                query = strategy.get("search_query", "")
                priority = strategy.get("priority", "medium")
                priority_icon = {"high": "🔴", "medium": "🟡", "low": "🟢"}.get(
                    priority, "🟡"
                )
                lines.append(f"### 6.{i} {priority_icon} {topic_name}\n")
                if brief:
                    lines.append(f"{brief}\n")
                if query:
                    lines.append(f"> 🔍 搜索建议：`{query}`\n")
        elif supplementary:
            lines.append("## 六、前置知识补充\n")
            lines.append("以下为理解本主题需要掌握的前置知识点：\n")
            for i, topic_item in enumerate(supplementary, 1):
                lines.append(f"### 6.{i} {topic_item}\n")
                lines.append("（建议通过搜索引擎进一步学习）\n")

        lines.append("## 七、学习路径\n")
        lines.append("```")
        lines.append(f"入门 → 理解 {safe_topic} 基础概念 →")
        lines.append("进阶 → 掌握核心方法与实践 →")
        lines.append("深入 → 研究架构原理与设计权衡 →")
        lines.append("实战 → 在真实项目中应用与迭代 →")
        lines.append("精通 → 形成体系化认知与创新能力")
        lines.append("```\n")

        lines.append("---\n")
        lines.append("*本文档由 Harness Engine v2.0 + deepseek-v4-flash 自动生成*")

        return "\n".join(lines)

    def _extract_ppt_paths(self, user_input: str) -> List[str]:
        """从用户输入中提取 PPT 文件路径 - 支持精确路径、文件名、目录"""
        paths = []
        found_set = set()

        search_dirs = self._get_search_dirs(user_input)

        words = user_input.split()
        ppt_extensions = (".ppt", ".pptx")

        for word in words:
            clean = word.strip("\"'(),")
            if not clean:
                continue

            if clean.endswith(ppt_extensions):
                if os.path.isabs(clean) and os.path.exists(clean):
                    abs_path = os.path.abspath(clean)
                    if abs_path not in found_set:
                        found_set.add(abs_path)
                        paths.append(abs_path)
                else:
                    for search_dir in search_dirs:
                        candidate = os.path.join(search_dir, clean)
                        if os.path.exists(candidate):
                            abs_path = os.path.abspath(candidate)
                            if abs_path not in found_set:
                                found_set.add(abs_path)
                                paths.append(abs_path)
                            break
                    else:
                        for search_dir in search_dirs:
                            found = self._glob_find(search_dir, clean)
                            for fp in found:
                                if fp not in found_set:
                                    found_set.add(fp)
                                    paths.append(fp)

            elif os.path.isdir(clean):
                ppt_files = self._scan_directory(clean)
                for fp in ppt_files:
                    if fp not in found_set:
                        found_set.add(fp)
                        paths.append(fp)

        if not paths:
            for search_dir in search_dirs:
                ppt_files = self._scan_directory(search_dir)
                for fp in ppt_files:
                    if fp not in found_set:
                        found_set.add(fp)
                        paths.append(fp)
                if paths:
                    break

        return paths

    def _get_search_dirs(self, user_input: str) -> List[str]:
        """获取 PPT 文件的搜索目录"""
        dirs = []
        cwd = os.getcwd()

        data_dir = os.path.join(cwd, "..", "data")
        if os.path.isdir(data_dir):
            dirs.append(os.path.abspath(data_dir))

        parent_data_dir = os.path.join(cwd, "data")
        if os.path.isdir(parent_data_dir):
            dirs.append(os.path.abspath(parent_data_dir))

        current_dir = os.path.dirname(os.path.abspath(__file__))
        for _ in range(5):
            candidate_data = os.path.join(current_dir, "data")
            if os.path.isdir(candidate_data):
                abs_dir = os.path.abspath(candidate_data)
                if abs_dir not in dirs:
                    dirs.append(abs_dir)

        user_dirs = re.findall(r'["\']?([^\s"\']+\.pptx?)["\']?', user_input)
        for ud in user_dirs:
            ud_dir = os.path.dirname(ud)
            if ud_dir and os.path.isdir(ud_dir):
                abs_dir = os.path.abspath(ud_dir)
                if abs_dir not in dirs:
                    dirs.append(abs_dir)

        week13_dir = os.path.join(cwd, "..")
        if os.path.isdir(week13_dir):
            dirs.append(os.path.abspath(week13_dir))

        if cwd not in dirs:
            dirs.append(cwd)

        return dirs

    def _scan_directory(self, directory: str) -> List[str]:
        """扫描目录中的 PPT 文件"""
        ppt_files = []
        try:
            for root, dirs, files in os.walk(directory):
                for f in files:
                    if f.lower().endswith((".ppt", ".pptx")):
                        ppt_files.append(os.path.join(root, f))
        except PermissionError:
            pass
        return ppt_files

    def _glob_find(self, base_dir: str, filename: str) -> List[str]:
        """在目录中递归查找文件"""
        found = []
        try:
            for root, dirs, files in os.walk(base_dir):
                for f in files:
                    if f == filename:
                        found.append(os.path.join(root, f))
                if found:
                    break
        except PermissionError:
            pass
        return found

    def _handle_no_match(self, user_input: str) -> dict:
        """无匹配 Skill 时的处理"""
        return {
            "status": "no_match",
            "message": "未匹配到可用 Skill",
            "available_skills": [s["name"] for s in self.disclosure.skill_index],
            "hint": "请尝试提供 PPT 文件路径",
        }


class Harness:
    """
    Harness 主引擎

    架构: Gateway → Progressive Disclosure → Agent Node
    记忆: 四层记忆模型 (Working → Short-term → Long-term → Vector)
    披露: 渐进式披露 (常驻层 → 触发层 → 执行层)
    """

    def __init__(self, config_dir: str):
        self.config_dir = config_dir
        self.gateway = Gateway()
        self.memory = MemorySystem(config_dir)
        self.disclosure = ProgressiveDisclosure(config_dir)

        soul_path = os.path.join(config_dir, "SOUL.md")
        self.llm = create_llm_client(soul_path=soul_path)
        if self.llm.is_available:
            print(f"  ✅ LLM 已就绪: {self.llm.model}")
        else:
            print(f"  ⚠️  LLM 未配置 (DEEPSEEK_API_KEY 未设置)，将使用规则降级模式")

        self.agent = AgentNode(self.gateway, self.memory, self.disclosure, llm=self.llm)
        self._display_banner()

    def _display_banner(self):
        """显示 Harness 启动信息"""
        print("=" * 60)
        print("  Harness Engine v2.0 (LLM-Powered)")
        print("  渐进式披露架构 | Progressive Disclosure")
        print("=" * 60)
        print(f"\n  配置目录: {self.config_dir}")
        print(f"  结果输出: {self.agent.result_dir}")
        llm_status = "✅ deepseek-v4-flash" if self.llm.is_available else "⚠️  降级模式"
        print(f"  LLM 引擎: {llm_status}")
        print(f"  可用 Skills: {len(self.disclosure.skill_index)}")
        for s in self.disclosure.skill_index:
            triggers = ", ".join(s.get("triggers", [])[:2])
            print(f"    • {s['name']} (触发: {triggers})")
        print()

    def process_input(self, user_input: str, session_id: str = None) -> dict:
        """处理用户输入"""
        if not session_id:
            session_id = f"session_{int(time.time())}"

        self.gateway.submit_message(session_id, user_input)
        result = self.agent.run(session_id, user_input)
        self.memory.flush()

        return result

    def get_status(self) -> dict:
        """获取 Harness 状态"""
        return {
            "sessions": len(self.gateway.sessions),
            "skills_loaded": list(self.disclosure.loaded_skills.keys()),
            "working_memory_size": len(self.memory.working_memory),
            "available_skills": [s["name"] for s in self.disclosure.skill_index],
            "result_dir": str(self.agent.result_dir),
        }

    def interactive_mode(self):
        """交互式模式"""
        print("\n进入交互模式，输入 'exit' 退出\n")
        while True:
            try:
                user_input = input("你: ").strip()
                if user_input.lower() in ["exit", "quit", "q", "退出"]:
                    print("\n[Harness] 退出交互模式")
                    break
                if not user_input:
                    continue

                result = self.process_input(user_input)
                self._display_result(result)

            except KeyboardInterrupt:
                print("\n\n[Harness] 中断，退出交互模式")
                break
            except Exception as e:
                print(f"\n[Harness] 错误: {e}")

    def _display_result(self, result: dict):
        """显示执行结果"""
        status = result.get("status", "unknown")
        print(f"\n[Result] 状态: {status}")

        if status == "success":
            print(f"  Skill: {result.get('skill', 'N/A')}")
            print(f"  完成步骤: {result.get('steps_completed', 0)}")
            for r in result.get("results", []):
                tool_result = r.get("result", {})
                if isinstance(tool_result, dict):
                    status_icon = "✓" if tool_result.get("status") == "success" else "✗"
                    extra = ""
                    if tool_result.get("output_path"):
                        extra = f" → {tool_result['output_path']}"
                    elif tool_result.get("file_size"):
                        extra = f" ({tool_result['file_size']} chars)"
                    print(
                        f"    [{status_icon}] {r['step']}: {r.get('tool', 'N/A')}{extra}"
                    )
                else:
                    print(f"    - {r['step']}: {r.get('tool', 'N/A')}")
            if result.get("output_path"):
                print(f"\n  📄 生成文件: {result['output_path']}")
        elif status == "no_match":
            print(f"  可用 Skills: {result.get('available_skills', [])}")
            print(f"  提示: {result.get('hint', '')}")
        elif status == "error":
            print(f"  错误: {result.get('message', '')}")

        print()
